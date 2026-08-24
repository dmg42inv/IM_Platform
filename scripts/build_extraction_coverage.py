"""Extraction / read-coverage layer for the all-folder document graph sandbox.

This is step 2 after copy/hash coverage (see scripts/bootstrap_sandbox_coverage.py).
It reads a single investment folder's copied archive, attempts to extract text
from every file using the parsers available in the current environment, and
records an explicit terminal status for EVERY file so nothing is silently
dropped. Outputs are written to the sandbox only, never the live source.

Terminal statuses (every file ends in exactly one):
    extracted_text          - real text captured
    extracted_empty         - parser ran but no meaningful text
    ocr_required            - image, or PDF with no text layer (scanned)
    parser_unavailable      - format needs a parser not installed here
    manual_review_required  - unknown / binary
    error                   - extraction raised, message recorded

Usage:
    python scripts/build_extraction_coverage.py --folder "<...>/_RS/AF/0_E_Q_U_I_T_Y/ONT/SR"
"""

from __future__ import annotations

import argparse
import csv
import json
import os
import re
from dataclasses import asdict, dataclass, field
from datetime import datetime, timezone
from pathlib import Path

TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".log"}
HTML_EXTENSIONS = {".html", ".htm", ".mhtml"}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
XLSX_EXTENSIONS = {".xlsx", ".xlsm"}
XLS_EXTENSIONS = {".xls"}
PPTX_EXTENSIONS = {".pptx"}
MSG_EXTENSIONS = {".msg"}
EML_EXTENSIONS = {".eml"}
IMAGE_EXTENSIONS = {".bmp", ".png", ".jpg", ".jpeg", ".jfif", ".jpe", ".tif", ".tiff", ".gif"}
DOC_EXTENSIONS = {".doc"}
RTF_EXTENSIONS = {".rtf"}
LEGACY_OFFICE_EXTENSIONS = {".ppt", ".pps"}
ARCHIVE_EXTENSIONS = {".zip"}

# A PDF whose extracted text averages fewer than this many characters per page
# is treated as scanned / image-only and routed to OCR instead of being called
# "extracted".
PDF_MIN_CHARS_PER_PAGE = 12
PDF_MIN_TOTAL_CHARS = 40


@dataclass
class ExtractionRecord:
    relative_path: str
    extension: str
    size_bytes: int
    parser: str
    status: str
    char_count: int = 0
    page_sheet_slide_count: int = 0
    ocr_required: bool = False
    manual_review_required: bool = False
    error_message: str = ""
    note: str = ""
    text_file: str = ""


def long_path(path: Path) -> str:
    resolved = str(path.resolve(strict=False))
    if resolved.startswith("\\\\?\\"):
        return resolved
    return "\\\\?\\" + resolved


def slugify(text: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9]+", "_", text).strip("_")
    return slug[:80] or "archive"


def _read_bytes(path: Path) -> bytes:
    with open(long_path(path), "rb") as handle:
        return handle.read()


def extract_text_like(path: Path) -> tuple[str, int]:
    raw = _read_bytes(path)
    text = raw.decode("utf-8", errors="replace")
    return text, 0


def extract_html(path: Path) -> tuple[str, int]:
    raw = _read_bytes(path)
    html = raw.decode("utf-8", errors="replace")
    html = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", html)
    text = re.sub(r"(?s)<[^>]+>", " ", html)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
    return text.strip(), 0


def extract_pdf(path: Path) -> tuple[str, int, bool]:
    """Return (text, page_count, ocr_required)."""
    import pdfplumber

    pages_text: list[str] = []
    page_count = 0
    with pdfplumber.open(long_path(path)) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            pages_text.append(page.extract_text() or "")
    text = "\n\n".join(pages_text).strip()
    avg = (len(text) / page_count) if page_count else 0
    scanned = len(text) < PDF_MIN_TOTAL_CHARS or avg < PDF_MIN_CHARS_PER_PAGE
    return text, page_count, scanned


def extract_pdf_via_pdfium(path: Path) -> tuple[str, int]:
    """Fallback PDF text extractor for files pdfminer/pdfplumber cannot parse."""
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(long_path(path))
    parts: list[str] = []
    page_count = len(pdf)
    try:
        for index in range(page_count):
            textpage = pdf[index].get_textpage()
            parts.append(textpage.get_text_range() or "")
    finally:
        pdf.close()
    return "\n\n".join(parts).strip(), page_count


def extract_eml(path: Path) -> tuple[str, int]:
    import email
    from email import policy

    with open(long_path(path), "rb") as handle:
        message = email.message_from_binary_file(handle, policy=policy.default)
    header = f"From: {message['from']}\nTo: {message['to']}\nDate: {message['date']}\nSubject: {message['subject']}\n\n"
    body = ""
    if message.is_multipart():
        for part in message.walk():
            if part.get_content_type() == "text/plain":
                body += part.get_content()
        if not body.strip():
            for part in message.walk():
                if part.get_content_type() == "text/html":
                    body += re.sub(r"(?s)<[^>]+>", " ", part.get_content())
    else:
        body = message.get_content() if message.get_content_type().startswith("text") else ""
    return (header + body).strip(), 0


def extract_docx(path: Path) -> tuple[str, int]:
    import docx

    document = docx.Document(long_path(path))
    parts: list[str] = [p.text for p in document.paragraphs if p.text and p.text.strip()]
    block_count = len(document.paragraphs)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
        block_count += 1
    return "\n".join(parts).strip(), block_count


def extract_xlsx(path: Path) -> tuple[str, int]:
    import openpyxl

    workbook = openpyxl.load_workbook(long_path(path), read_only=True, data_only=True)
    lines: list[str] = []
    sheet_count = 0
    try:
        for sheet in workbook.worksheets:
            sheet_count += 1
            try:
                dims = sheet.calculate_dimension()
            except Exception:
                dims = "unknown"
            lines.append(f"# Sheet: {sheet.title} (dims={dims})")
            try:
                for row in sheet.iter_rows(values_only=True):
                    values = [str(v) for v in row if v is not None and str(v).strip()]
                    if values:
                        lines.append("\t".join(values))
            except Exception as exc:  # noqa: BLE001 - keep other sheets
                lines.append(f"# (row read error: {exc.__class__.__name__}: {exc})")
    finally:
        workbook.close()
    return "\n".join(lines).strip(), sheet_count


def extract_xls(path: Path) -> tuple[str, int]:
    import xlrd

    book = xlrd.open_workbook(long_path(path))
    lines: list[str] = []
    for sheet in book.sheets():
        lines.append(f"# Sheet: {sheet.name} ({sheet.nrows}x{sheet.ncols})")
        for r in range(sheet.nrows):
            values = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
            values = [v for v in values if v.strip()]
            if values:
                lines.append("\t".join(values))
    return "\n".join(lines).strip(), book.nsheets


def extract_doc_via_word(path: Path) -> tuple[str, int]:
    """Read a legacy binary .doc via Word COM automation (Windows + Word only)."""
    import pythoncom
    import win32com.client as win32

    pythoncom.CoInitialize()
    word = win32.DispatchEx("Word.Application")
    word.Visible = False
    word.DisplayAlerts = 0
    try:
        document = word.Documents.Open(
            str(path.resolve()),
            ReadOnly=True,
            ConfirmConversions=False,
            AddToRecentFiles=False,
        )
        try:
            text = document.Content.Text or ""
        finally:
            document.Close(False)
        return text.strip(), 0
    finally:
        word.Quit()
        pythoncom.CoUninitialize()


def _ensure_tesseract() -> bool:
    """Point pytesseract at an installed tesseract.exe if it is not on PATH."""
    import pytesseract

    candidates = [
        os.environ.get("TESSERACT_CMD"),
        os.path.join(os.environ.get("ProgramFiles", ""), "Tesseract-OCR", "tesseract.exe"),
        os.path.join(os.environ.get("ProgramFiles(x86)", ""), "Tesseract-OCR", "tesseract.exe"),
        os.path.join(os.environ.get("LOCALAPPDATA", ""), "Programs", "Tesseract-OCR", "tesseract.exe"),
    ]
    for candidate in candidates:
        if candidate and os.path.exists(candidate):
            pytesseract.pytesseract.tesseract_cmd = candidate
            return True
    try:
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def try_ocr(path: Path) -> tuple[str, bool]:
    """Attempt OCR of an image. Return (text, ocr_succeeded). Requires tesseract binary."""
    try:
        import pytesseract
        from PIL import Image

        if not _ensure_tesseract():
            return "", False
        with Image.open(long_path(path)) as img:
            text = pytesseract.image_to_string(img)
        return text.strip(), True
    except Exception:
        return "", False


def ocr_pdf(path: Path) -> tuple[str, int]:
    """Render each PDF page to an image and OCR it. Requires tesseract + pypdfium2."""
    import pypdfium2 as pdfium
    import pytesseract

    if not _ensure_tesseract():
        raise RuntimeError("tesseract binary not available")
    pdf = pdfium.PdfDocument(long_path(path))
    parts: list[str] = []
    page_count = len(pdf)
    try:
        for index in range(page_count):
            page = pdf[index]
            bitmap = page.render(scale=2.0)
            image = bitmap.to_pil()
            parts.append(pytesseract.image_to_string(image))
    finally:
        pdf.close()
    return "\n\n".join(parts).strip(), page_count


def route_and_extract(path: Path, ext: str, do_ocr: bool) -> tuple[str, ExtractionRecord]:
    size = path.stat().st_size
    rec = ExtractionRecord(
        relative_path="",  # filled by caller
        extension=ext,
        size_bytes=size,
        parser="",
        status="",
    )
    text = ""

    try:
        if ext in TEXT_EXTENSIONS:
            rec.parser = "text_decode"
            text, _ = extract_text_like(path)
        elif ext in HTML_EXTENSIONS:
            rec.parser = "html_striptags"
            text, _ = extract_html(path)
        elif ext in PDF_EXTENSIONS:
            rec.parser = "pdfplumber"
            try:
                text, pages, scanned = extract_pdf(path)
                rec.page_sheet_slide_count = pages
            except Exception as pdf_exc:  # noqa: BLE001 - malformed PDF; try fallbacks
                text, pages = "", 0
                try:
                    text, pages = extract_pdf_via_pdfium(path)
                    rec.parser = "pypdfium2"
                    rec.page_sheet_slide_count = pages
                except Exception:  # noqa: BLE001
                    text = ""
                if not text.strip():
                    if do_ocr:
                        try:
                            ocr_text, _ = ocr_pdf(path)
                            if ocr_text.strip():
                                rec.parser = "pypdfium2+ocr"
                                rec.status = "extracted_text"
                                rec.char_count = len(ocr_text)
                                return ocr_text, rec
                        except Exception:  # noqa: BLE001
                            pass
                    rec.status = "manual_review_required"
                    rec.manual_review_required = True
                    rec.note = f"PDF parse failed ({pdf_exc.__class__.__name__}); fallback produced no text."
                    return "", rec
                scanned = False
            if scanned:
                if do_ocr:
                    try:
                        ocr_text, _ = ocr_pdf(path)
                    except Exception as exc:  # noqa: BLE001
                        rec.ocr_required = True
                        rec.status = "ocr_required"
                        rec.note = f"Scanned PDF; OCR failed: {exc.__class__.__name__}: {exc}"
                        return "", rec
                    if ocr_text.strip():
                        rec.parser = "pdfplumber+ocr"
                        rec.status = "extracted_text"
                        rec.char_count = len(ocr_text)
                        return ocr_text, rec
                    rec.ocr_required = True
                    rec.status = "ocr_required"
                    rec.note = "Scanned PDF; OCR produced no text."
                    return "", rec
                rec.ocr_required = True
                rec.status = "ocr_required"
                rec.note = "PDF has little/no text layer (likely scanned)."
                return text, rec
        elif ext in DOCX_EXTENSIONS:
            rec.parser = "python-docx"
            try:
                text, blocks = extract_docx(path)
                rec.page_sheet_slide_count = blocks
            except Exception as exc:  # noqa: BLE001
                if exc.__class__.__name__ == "PackageNotFoundError":
                    rec.status = "manual_review_required"
                    rec.manual_review_required = True
                    rec.note = "Not a valid OOXML package (possibly legacy .doc renamed or corrupt); manual review."
                    return "", rec
                raise
        elif ext in XLSX_EXTENSIONS:
            rec.parser = "openpyxl"
            try:
                text, sheets = extract_xlsx(path)
                rec.page_sheet_slide_count = sheets
            except Exception as xlsx_exc:  # noqa: BLE001 - not-a-zip / corrupt; try xlrd
                try:
                    text, sheets = extract_xls(path)
                    rec.parser = "xlrd_fallback"
                    rec.page_sheet_slide_count = sheets
                except Exception as xls_exc:  # noqa: BLE001
                    encrypted = "encrypt" in f"{xlsx_exc}{xls_exc}".lower()
                    rec.status = "manual_review_required"
                    rec.manual_review_required = True
                    rec.note = (
                        "Encrypted/password-protected workbook; cannot read."
                        if encrypted
                        else f"Workbook unreadable ({xlsx_exc.__class__.__name__}/{xls_exc.__class__.__name__}); manual review."
                    )
                    return "", rec
        elif ext in XLS_EXTENSIONS:
            rec.parser = "xlrd"
            text, sheets = extract_xls(path)
            rec.page_sheet_slide_count = sheets
        elif ext in IMAGE_EXTENSIONS:
            rec.parser = "ocr"
            if do_ocr:
                text, ok = try_ocr(path)
                if ok and text:
                    rec.status = "extracted_text"
                    rec.char_count = len(text)
                    return text, rec
            rec.ocr_required = True
            rec.status = "ocr_required"
            rec.note = "Image file; OCR required."
            return text, rec
        elif ext in PPTX_EXTENSIONS:
            rec.parser = "python-pptx"
            try:
                from pptx import Presentation  # type: ignore

                prs = Presentation(long_path(path))
                parts: list[str] = []
                slide_count = 0
                for slide in prs.slides:
                    slide_count += 1
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                line = "".join(run.text for run in para.runs)
                                if line.strip():
                                    parts.append(line)
                text = "\n".join(parts).strip()
                rec.page_sheet_slide_count = slide_count
            except ModuleNotFoundError:
                rec.status = "parser_unavailable"
                rec.manual_review_required = True
                rec.note = "python-pptx not installed."
                return "", rec
        elif ext in MSG_EXTENSIONS:
            rec.parser = "extract-msg"
            try:
                import extract_msg  # type: ignore

                msg = extract_msg.Message(long_path(path))
                header = f"From: {msg.sender}\nTo: {msg.to}\nDate: {msg.date}\nSubject: {msg.subject}\n\n"
                text = (header + (msg.body or "")).strip()
            except ModuleNotFoundError:
                rec.status = "parser_unavailable"
                rec.manual_review_required = True
                rec.note = "extract-msg not installed."
                return "", rec
        elif ext in EML_EXTENSIONS:
            rec.parser = "email_stdlib"
            text, _ = extract_eml(path)
        elif ext in DOC_EXTENSIONS or ext in RTF_EXTENSIONS:
            rec.parser = "word_com"
            try:
                text, _ = extract_doc_via_word(path)
            except Exception as exc:  # noqa: BLE001 - Word/pywin32 may be absent or COM may fail
                rec.status = "parser_unavailable"
                rec.manual_review_required = True
                rec.note = f"Legacy .doc; Word COM conversion unavailable/failed: {exc.__class__.__name__}: {exc}"
                return "", rec
        elif ext in LEGACY_OFFICE_EXTENSIONS:
            rec.parser = "none"
            rec.status = "parser_unavailable"
            rec.manual_review_required = True
            rec.note = "Legacy binary Office format; convert to modern format or parse manually."
            return "", rec
        else:
            rec.parser = "none"
            rec.status = "manual_review_required"
            rec.manual_review_required = True
            rec.note = "Unknown/binary file type."
            return "", rec
    except Exception as exc:  # noqa: BLE001 - we must record, not crash the run
        rec.status = "error"
        rec.error_message = f"{exc.__class__.__name__}: {exc}"
        return "", rec

    # Common finalisation for parsers that produced text.
    rec.char_count = len(text)
    if text.strip():
        rec.status = "extracted_text"
    else:
        rec.status = "extracted_empty"
        rec.note = rec.note or "Parser ran but produced no text."
    return text, rec


def process_folder(folder: Path, do_ocr: bool) -> dict[str, object]:
    index_root = folder / "00_Index"
    inventory_path = index_root / "file_readiness_inventory.csv"
    archive_root = folder / "99_Archive"
    intel_root = index_root / "Document_Intelligence"
    text_root = intel_root / "extracted_text"

    if not inventory_path.exists():
        raise FileNotFoundError(f"Missing readiness inventory: {inventory_path}")
    if not archive_root.exists():
        raise FileNotFoundError(f"Missing archive copy: {archive_root}")

    os.makedirs(long_path(intel_root), exist_ok=True)
    os.makedirs(long_path(text_root), exist_ok=True)

    with inventory_path.open("r", encoding="utf-8", newline="") as handle:
        inventory = list(csv.DictReader(handle))

    records: list[ExtractionRecord] = []
    for row in inventory:
        relative = row["relative_path"]
        ext = (row.get("extension") or Path(relative).suffix).lower()
        source_file = archive_root / relative
        if not source_file.exists():
            rec = ExtractionRecord(
                relative_path=relative,
                extension=ext,
                size_bytes=0,
                parser="none",
                status="error",
                error_message="File listed in inventory but missing from archive copy.",
            )
            records.append(rec)
            continue

        text, rec = route_and_extract(source_file, ext, do_ocr)
        rec.relative_path = relative

        if text.strip():
            text_target = text_root / (relative + ".txt")
            os.makedirs(long_path(text_target.parent), exist_ok=True)
            with open(long_path(text_target), "w", encoding="utf-8") as out:
                out.write(text)
            rec.text_file = str(text_target.relative_to(intel_root))
        records.append(rec)

        if ext in ARCHIVE_EXTENSIONS:
            staging_root = folder / "98_Unzipped"
            members, error = expand_zip(source_file, relative, staging_root, do_ocr)
            rec.status = "archive_expanded"
            rec.parser = "zip"
            rec.note = error or f"Archive expanded ({len(members)} members)."
            for member_rel, member_text, member_rec in members:
                _write_text_file(text_root, intel_root, member_rel, member_text, member_rec)
                records.append(member_rec)

    return _write_outputs(folder, intel_root, records, do_ocr)


def _write_text_file(text_root: Path, intel_root: Path, relative: str, text: str, rec: ExtractionRecord) -> None:
    if text.strip():
        text_target = text_root / (relative + ".txt")
        os.makedirs(long_path(text_target.parent), exist_ok=True)
        with open(long_path(text_target), "w", encoding="utf-8") as out:
            out.write(text)
        rec.text_file = str(text_target.relative_to(intel_root))


def expand_zip(source_zip: Path, relative: str, staging_root: Path, do_ocr: bool, depth: int = 0) -> tuple[list[tuple[str, str, ExtractionRecord]], str]:
    """Extract a zip's members and run extraction on each. Returns (member_results, error)."""
    import zipfile

    results: list[tuple[str, str, ExtractionRecord]] = []
    dest = staging_root / slugify(relative)
    try:
        with zipfile.ZipFile(long_path(source_zip)) as archive:
            members = [name for name in archive.namelist() if not name.endswith("/")]
            for name in members:
                base = name.replace("\\", "/").split("/")[-1]
                if "__MACOSX" in name or base == ".DS_Store" or base.startswith("._"):
                    continue  # macOS archive metadata, not a real document
                safe_member = name.replace("\\", "/")
                safe_member = "/".join(part.replace(":", "_") for part in safe_member.split("/"))
                target = dest / safe_member
                os.makedirs(long_path(target.parent), exist_ok=True)
                with open(long_path(target), "wb") as out:
                    out.write(archive.read(name))
                member_ext = Path(safe_member).suffix.lower()
                member_rel = f"{relative}__unzipped/{safe_member}"
                if member_ext in ARCHIVE_EXTENSIONS and depth < 1:
                    nested, nested_err = expand_zip(target, member_rel, staging_root, do_ocr, depth + 1)
                    results.extend(nested)
                    container = ExtractionRecord(
                        relative_path=member_rel, extension=member_ext,
                        size_bytes=target.stat().st_size, parser="zip",
                        status="archive_expanded",
                        note=(nested_err or f"Nested archive expanded ({len(nested)} members)."),
                    )
                    results.append((member_rel, "", container))
                else:
                    member_text, member_rec = route_and_extract(target, member_ext, do_ocr)
                    member_rec.relative_path = member_rel
                    results.append((member_rel, member_text, member_rec))
    except Exception as exc:  # noqa: BLE001
        return results, f"{exc.__class__.__name__}: {exc}"
    return results, ""


def _write_outputs(folder: Path, intel_root: Path, records: list[ExtractionRecord], do_ocr: bool) -> dict[str, object]:
    status_csv = intel_root / "extraction_status.csv"
    fieldnames = list(asdict(records[0]).keys()) if records else ["relative_path"]
    with open(long_path(status_csv), "w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames)
        writer.writeheader()
        for rec in records:
            writer.writerow(asdict(rec))

    status_counts: dict[str, int] = {}
    parser_counts: dict[str, int] = {}
    total_chars = 0
    for rec in records:
        status_counts[rec.status] = status_counts.get(rec.status, 0) + 1
        parser_counts[rec.parser] = parser_counts.get(rec.parser, 0) + 1
        total_chars += rec.char_count

    extracted = status_counts.get("extracted_text", 0)
    blocked = sum(
        status_counts.get(s, 0)
        for s in ("ocr_required", "parser_unavailable", "manual_review_required", "error")
    )

    manifest = {
        "folder": str(folder),
        "generated_at_utc": datetime.now(timezone.utc).isoformat(),
        "ocr_attempted": do_ocr,
        "file_count": len(records),
        "extracted_text_count": extracted,
        "extracted_empty_count": status_counts.get("extracted_empty", 0),
        "blocked_count": blocked,
        "total_extracted_chars": total_chars,
        "status_counts": status_counts,
        "parser_counts": parser_counts,
    }
    with open(long_path(intel_root / "extraction_manifest.json"), "w", encoding="utf-8") as handle:
        handle.write(json.dumps(manifest, indent=2))

    with open(long_path(intel_root / "EXTRACTION_STATUS.md"), "w", encoding="utf-8") as handle:
        handle.write(_render_status_md(folder, manifest, status_counts))

    return manifest


def _render_status_md(folder: Path, manifest: dict[str, object], status_counts: dict[str, int]) -> str:
    lines = [
        f"# Extraction Coverage - {folder.parent.name}",
        "",
        f"- Generated (UTC): {manifest['generated_at_utc']}",
        f"- Files in inventory: {manifest['file_count']}",
        f"- Extracted with text: {manifest['extracted_text_count']}",
        f"- Extracted but empty: {manifest['extracted_empty_count']}",
        f"- Blocked (OCR / parser / manual / error): {manifest['blocked_count']}",
        f"- Total extracted characters: {manifest['total_extracted_chars']}",
        "",
        "## Status breakdown",
        "",
    ]
    for status, count in sorted(status_counts.items(), key=lambda kv: (-kv[1], kv[0])):
        lines.append(f"- {status}: {count}")
    lines.extend([
        "",
        "## Interpretation",
        "",
        "Every file in file_readiness_inventory.csv has exactly one terminal status here.",
        "extracted_text files have real captured text under Document_Intelligence/extracted_text.",
        "Blocked files are NOT graph-ready yet and must be resolved (install parser, run OCR, or manual review)",
        "before this folder can be called graph-complete.",
    ])
    return "\n".join(lines) + "\n"


def main() -> int:
    parser = argparse.ArgumentParser(description="Extraction/read-coverage layer for one investment sandbox folder.")
    parser.add_argument("--folder", type=Path, required=True, help="Path to an investment folder's SR root (contains 00_Index and 99_Archive).")
    parser.add_argument("--ocr", action="store_true", help="Attempt OCR for images if a tesseract binary is available.")
    args = parser.parse_args()

    manifest = process_folder(args.folder, args.ocr)
    print(json.dumps(manifest, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
