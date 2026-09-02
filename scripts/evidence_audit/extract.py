"""Stage 5 of the evidence coverage audit: extract text into the truth layer, keyed on the hash.

Every node records where in the document it came from - page number for PDFs, slide for decks,
sheet and row span for workbooks - because a citation that cannot name a page is not a citation.

Documents are keyed on SHA-256 rather than on a path or a filename, so a file that is copied,
renamed or moved is still recognised as the same evidence, and two different files that happen
to share a name are never conflated.

Re-running is safe: a document already extracted at the same hash is skipped unless --force.

    .\\.venv\\Scripts\\python.exe -m scripts.evidence_audit.extract --folder "4. MGX"
    .\\.venv\\Scripts\\python.exe -m scripts.evidence_audit.extract --all
"""

from __future__ import annotations

import argparse
import sqlite3
import sys
import time
import traceback
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[2]
AUDIT_DB = REPO_ROOT / "data" / "evidence" / "audit.sqlite"

MAX_NODE_CHARS = 2400

SCHEMA = """
CREATE TABLE IF NOT EXISTS documents (
    sha256          TEXT PRIMARY KEY,
    source_id       TEXT,
    full_path       TEXT,
    filename        TEXT,
    extension       TEXT,
    extractor       TEXT,
    extraction_status TEXT NOT NULL,
    error           TEXT,
    page_count      INTEGER,
    node_count      INTEGER,
    char_count      INTEGER,
    extracted_at    TEXT
);
CREATE TABLE IF NOT EXISTS nodes (
    node_id     INTEGER PRIMARY KEY AUTOINCREMENT,
    sha256      TEXT NOT NULL,
    ordinal     INTEGER NOT NULL,
    locator     TEXT NOT NULL,
    page_no     INTEGER,
    kind        TEXT,
    heading     TEXT,
    text        TEXT NOT NULL
);
CREATE INDEX IF NOT EXISTS ix_nodes_sha ON nodes(sha256);
CREATE INDEX IF NOT EXISTS ix_docs_status ON documents(extraction_status);
"""


def long_path(path: str) -> str:
    return path if path.startswith("\\\\?\\") else "\\\\?\\" + path


def chunk(text: str, limit: int = MAX_NODE_CHARS) -> list[str]:
    """Split on blank lines, keeping paragraphs whole where they fit."""
    text = text.strip()
    if len(text) <= limit:
        return [text] if text else []
    out, current = [], ""
    for para in text.split("\n\n"):
        if len(current) + len(para) + 2 > limit and current:
            out.append(current.strip())
            current = ""
        if len(para) > limit:
            for i in range(0, len(para), limit):
                out.append(para[i:i + limit].strip())
        else:
            current += para + "\n\n"
    if current.strip():
        out.append(current.strip())
    return [o for o in out if o]


def extract_pdf(path: str) -> tuple[list[dict], int, str]:
    import pymupdf
    doc = pymupdf.open(long_path(path))
    nodes: list[dict] = []
    try:
        for page_index, page in enumerate(doc, start=1):
            for part, body in enumerate(chunk(page.get_text() or ""), start=1):
                nodes.append({"locator": f"page {page_index}" + (f" ({part})" if part > 1 else ""),
                              "page_no": page_index, "kind": "page", "heading": "", "text": body})
        pages = len(doc)
    finally:
        doc.close()
    return nodes, pages, "pymupdf"


def extract_docx(path: str) -> tuple[list[dict], int, str]:
    import docx
    document = docx.Document(long_path(path))
    nodes: list[dict] = []
    heading, buffer, block = "", [], 1

    def flush():
        nonlocal buffer, block
        joined = "\n\n".join(buffer).strip()
        for body in chunk(joined):
            nodes.append({"locator": f"block {block}", "page_no": None, "kind": "paragraph",
                          "heading": heading, "text": body})
            block += 1
        buffer = []

    for para in document.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        if (para.style.name or "").lower().startswith("heading"):
            flush()
            heading = text
        buffer.append(text)
        if sum(len(b) for b in buffer) > MAX_NODE_CHARS:
            flush()
    flush()
    for table_index, table in enumerate(document.tables, start=1):
        rows = [" | ".join(c.text.strip() for c in row.cells) for row in table.rows]
        for body in chunk("\n".join(rows)):
            nodes.append({"locator": f"table {table_index}", "page_no": None, "kind": "table",
                          "heading": heading, "text": body})
    return nodes, 0, "python-docx"


def extract_xlsx(path: str) -> tuple[list[dict], int, str]:
    import openpyxl
    book = openpyxl.load_workbook(long_path(path), data_only=True, read_only=True)
    nodes: list[dict] = []
    try:
        for sheet in book.worksheets:
            lines, first_row = [], None
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                cells = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if not cells:
                    continue
                first_row = first_row or row_index
                lines.append(f"r{row_index}: " + " | ".join(cells))
                if sum(len(x) for x in lines) > MAX_NODE_CHARS:
                    nodes.append({"locator": f"{sheet.title}!r{first_row}-r{row_index}",
                                  "page_no": None, "kind": "sheet", "heading": sheet.title,
                                  "text": "\n".join(lines)})
                    lines, first_row = [], None
            if lines:
                nodes.append({"locator": f"{sheet.title}!r{first_row}-r{row_index}",
                              "page_no": None, "kind": "sheet", "heading": sheet.title,
                              "text": "\n".join(lines)})
    finally:
        book.close()
    return nodes, 0, "openpyxl"


def extract_pptx(path: str) -> tuple[list[dict], int, str]:
    from pptx import Presentation
    deck = Presentation(long_path(path))
    nodes: list[dict] = []
    for slide_index, slide in enumerate(deck.slides, start=1):
        pieces = [sh.text_frame.text.strip() for sh in slide.shapes
                  if sh.has_text_frame and sh.text_frame.text.strip()]
        title = pieces[0].splitlines()[0][:120] if pieces else ""
        for body in chunk("\n\n".join(pieces)):
            nodes.append({"locator": f"slide {slide_index}", "page_no": slide_index,
                          "kind": "slide", "heading": title, "text": body})
    return nodes, len(deck.slides), "python-pptx"


def extract_msg(path: str) -> tuple[list[dict], int, str]:
    import extract_msg
    message = extract_msg.Message(long_path(path))
    header = (f"From: {message.sender}\nTo: {message.to}\nDate: {message.date}\n"
              f"Subject: {message.subject}")
    nodes = [{"locator": f"message ({i})" if i > 1 else "message", "page_no": None,
              "kind": "email", "heading": message.subject or "", "text": body}
             for i, body in enumerate(chunk(header + "\n\n" + (message.body or "")), start=1)]
    message.close()
    return nodes, 0, "extract-msg"


def extract_text(path: str) -> tuple[list[dict], int, str]:
    with open(long_path(path), "r", encoding="utf-8", errors="replace") as handle:
        content = handle.read()
    nodes = [{"locator": f"block {i}", "page_no": None, "kind": "text", "heading": "", "text": body}
             for i, body in enumerate(chunk(content), start=1)]
    return nodes, 0, "plain-text"


EXTRACTORS = {
    ".pdf": extract_pdf, ".docx": extract_docx, ".xlsx": extract_xlsx, ".xlsm": extract_xlsx,
    ".pptx": extract_pptx, ".msg": extract_msg, ".txt": extract_text, ".csv": extract_text,
}
UNSUPPORTED = {".doc": "legacy Word binary; needs conversion to .docx",
               ".xls": "legacy Excel binary; needs conversion to .xlsx",
               ".ppt": "legacy PowerPoint binary; needs conversion to .pptx",
               ".eml": "no .eml parser installed"}


def run(conn: sqlite3.Connection, folder: str | None, limit: int | None, force: bool,
        progress_every: int) -> None:
    where = "hash_status='hashed' AND sha256 IS NOT NULL"
    params: list = []
    if folder:
        where += " AND full_path LIKE ?"
        params.append(f"%{folder}%")
    rows = conn.execute(
        f"SELECT sha256, source_id, full_path, filename, extension, MIN(rowid) FROM sources"
        f" WHERE {where} GROUP BY sha256 ORDER BY full_path", params).fetchall()
    if limit:
        rows = rows[:limit]

    done = {r[0] for r in conn.execute(
        "SELECT sha256 FROM documents WHERE extraction_status='extracted'")} if not force else set()

    print(f"  {len(rows):,} unique documents in scope; {len(done):,} already extracted", flush=True)
    started = time.time()
    extracted = skipped = failed = unsupported = 0
    total_nodes = 0

    for index, (sha, source_id, full_path, filename, extension, _rid) in enumerate(rows, start=1):
        if sha in done:
            skipped += 1
            continue
        handler = EXTRACTORS.get(extension)
        if handler is None:
            conn.execute(
                "INSERT OR REPLACE INTO documents (sha256, source_id, full_path, filename,"
                " extension, extractor, extraction_status, error, page_count, node_count,"
                " char_count, extracted_at) VALUES (?,?,?,?,?,?,?,?,0,0,0,datetime('now'))",
                (sha, source_id, full_path, filename, extension, "none", "unsupported",
                 UNSUPPORTED.get(extension, "no extractor for this type")))
            unsupported += 1
            continue
        try:
            nodes, pages, extractor = handler(full_path)
            conn.execute("DELETE FROM nodes WHERE sha256=?", (sha,))
            conn.executemany(
                "INSERT INTO nodes (sha256, ordinal, locator, page_no, kind, heading, text)"
                " VALUES (?,?,?,?,?,?,?)",
                [(sha, i, n["locator"], n["page_no"], n["kind"], n["heading"], n["text"])
                 for i, n in enumerate(nodes, start=1)])
            chars = sum(len(n["text"]) for n in nodes)
            status = "extracted" if nodes else "extracted_empty"
            conn.execute(
                "INSERT OR REPLACE INTO documents (sha256, source_id, full_path, filename,"
                " extension, extractor, extraction_status, error, page_count, node_count,"
                " char_count, extracted_at) VALUES (?,?,?,?,?,?,?,NULL,?,?,?,datetime('now'))",
                (sha, source_id, full_path, filename, extension, extractor, status,
                 pages, len(nodes), chars))
            extracted += 1
            total_nodes += len(nodes)
        except Exception as exc:  # noqa: BLE001 - the reason must be recorded, not raised
            conn.execute(
                "INSERT OR REPLACE INTO documents (sha256, source_id, full_path, filename,"
                " extension, extractor, extraction_status, error, page_count, node_count,"
                " char_count, extracted_at) VALUES (?,?,?,?,?,?,?,?,0,0,0,datetime('now'))",
                (sha, source_id, full_path, filename, extension, "error", "error",
                 f"{type(exc).__name__}: {exc}"[:400]))
            failed += 1
        if index % progress_every == 0:
            conn.commit()
            elapsed = max(time.time() - started, 0.001)
            print(f"    {index:>6,}/{len(rows):,} | {extracted:>5,} extracted | "
                  f"{total_nodes:>7,} nodes | {failed:>3} failed | {unsupported:>3} unsupported | "
                  f"{index/elapsed:5.1f} docs/s", flush=True)
    conn.commit()

    print(f"\n  extracted {extracted:,} documents into {total_nodes:,} nodes "
          f"in {time.time()-started:,.0f}s")
    print(f"  {skipped:,} already done, {failed:,} failed, {unsupported:,} unsupported")
    for status, count in conn.execute(
            "SELECT extraction_status, COUNT(*) FROM documents GROUP BY 1 ORDER BY 2 DESC"):
        print(f"    {status:<18} {count:>7,}")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--db", type=Path, default=AUDIT_DB)
    parser.add_argument("--folder", type=str, default=None, help="Only paths containing this text")
    parser.add_argument("--all", action="store_true")
    parser.add_argument("--limit", type=int, default=None)
    parser.add_argument("--force", action="store_true")
    parser.add_argument("--progress-every", type=int, default=100)
    args = parser.parse_args()

    if not args.folder and not args.all:
        parser.error("pass --folder to scope the run, or --all")

    conn = sqlite3.connect(str(args.db))
    conn.executescript(SCHEMA)
    conn.execute("PRAGMA journal_mode=WAL")
    try:
        run(conn, args.folder, args.limit, args.force, args.progress_every)
    finally:
        conn.close()
    return 0


if __name__ == "__main__":
    sys.exit(main())
